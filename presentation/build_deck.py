#!/usr/bin/env python3
"""
PingIt SDK — Presentation generator
===================================

Builds `PingIt_SDK.pptx`, a fully-editable PowerPoint deck.

HOW TO EDIT
-----------
  1. Open the generated .pptx in PowerPoint / Keynote / Google Slides and edit
     text + move shapes directly. (Every text box and diagram box is native.)
  2. Edit the TEXT in this file (the `slide_*` functions in SECTION 3 — wording
     lives right there) or the THEME block, then regenerate:

        ./.venv/bin/python presentation/build_deck.py

STRUCTURE OF THIS FILE
----------------------
  SECTION 1 — THEME            colours, fonts, sizes (change the look here)
  SECTION 2 — HELPERS          small reusable builders
  SECTION 3 — SLIDES           one function per slide; the wording is inline
  SECTION 4 — BUILD            the running order of the slides

Source material: docs/PRD.md and docs/SERVER_ARCHITECTURE.md.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ======================================================================
# SECTION 1 — THEME  (edit colours / fonts here to restyle the whole deck)
# ======================================================================

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY    = RGBColor(0x0F, 0x2A, 0x43)
NAVY2   = RGBColor(0x16, 0x3A, 0x5C)
TEAL    = RGBColor(0x1F, 0xA8, 0xC9)
TEAL_DK = RGBColor(0x14, 0x7A, 0x92)
GREEN   = RGBColor(0x2E, 0xB8, 0x72)
RED     = RGBColor(0xE5, 0x5B, 0x5B)
AMBER   = RGBColor(0xE8, 0xA8, 0x3A)
LIGHT   = RGBColor(0xF4, 0xF7, 0xF9)
LIGHT2  = RGBColor(0xE6, 0xEE, 0xF3)
INK     = RGBColor(0x1A, 0x2B, 0x3C)
MUTED   = RGBColor(0x6B, 0x7C, 0x8D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PALE    = RGBColor(0xA9, 0xC0, 0xCE)
CODE_BG = RGBColor(0x10, 0x20, 0x30)
CODE_FG = RGBColor(0xCB, 0xE6, 0xF2)

FONT      = "Calibri"
FONT_MONO = "Consolas"

# ======================================================================
# SECTION 2 — HELPERS
# ======================================================================

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.level = ln.get("level", 0)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        p.space_before = Pt(ln.get("space_before", 0))
        if ln.get("line_spacing"):
            p.line_spacing = ln["line_spacing"]
        prefix = ""
        if ln.get("bullet"):
            prefix = "•  " if ln.get("level", 0) == 0 else "–  "
        run = p.add_run()
        run.text = prefix + ln["text"]
        f = run.font
        f.size = Pt(ln.get("size", 16))
        f.bold = ln.get("bold", False)
        f.italic = ln.get("italic", False)
        f.name = ln.get("font", FONT)
        f.color.rgb = ln.get("color", INK)
    return tb


def slide_base(title, kicker=None, accent=TEAL):
    s = prs.slides.add_slide(BLANK)
    rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
    _set_fill(rail, accent); _no_line(rail)
    if kicker:
        textbox(s, Inches(0.55), Inches(0.30), Inches(12.3), Inches(0.35),
                [{"text": kicker.upper(), "size": 13, "bold": True, "color": accent}])
    textbox(s, Inches(0.52), Inches(0.62), Inches(12.4), Inches(0.9),
            [{"text": title, "size": 29, "bold": True, "color": NAVY}])
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.46),
                             Inches(12.2), Pt(2))
    _set_fill(div, LIGHT2); _no_line(div)
    return s


def panel(slide, x, y, w, h, fill=LIGHT, line=None, line_w=1.0, radius=True):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    p = slide.shapes.add_shape(shp_type, x, y, w, h)
    _set_fill(p, fill)
    if line:
        _line(p, line, line_w)
    else:
        _no_line(p)
    try:
        p.adjustments[0] = 0.06
    except Exception:
        pass
    p.text_frame.word_wrap = True
    return p


def titled_panel(slide, x, y, w, h, title, color, body_lines, body_size=12,
                 head_h=Inches(0.55), tsize=13.5):
    """Panel with a coloured header strip and bullet body. body_lines = list of str."""
    panel(slide, x, y, w, h, fill=LIGHT, line=color, line_w=1.5)
    head = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, head_h)
    _set_fill(head, color); _no_line(head)
    tf = head.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = title
    rr.font.bold = True; rr.font.size = Pt(tsize); rr.font.color.rgb = WHITE; rr.font.name = FONT
    lines = []
    for b in body_lines:
        if isinstance(b, tuple):
            txt, mono = b
        else:
            txt, mono = b, False
        lines.append({"text": txt, "size": body_size, "color": INK, "bullet": True,
                      "space_after": 7, "font": (FONT_MONO if mono else FONT)})
    textbox(slide, x + Inches(0.22), y + head_h + Inches(0.12), w - Inches(0.42),
            h - head_h - Inches(0.2), lines)


def box(slide, x, y, w, h, title, sub=None, fill=NAVY, fg=WHITE,
        tsize=14, ssize=10.5, radius=True, line=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    b = slide.shapes.add_shape(shp_type, x, y, w, h)
    _set_fill(b, fill)
    if line:
        _line(b, line, 1.25)
    else:
        _no_line(b)
    try:
        b.adjustments[0] = 0.10
    except Exception:
        pass
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(tsize); r.font.bold = True; r.font.color.rgb = fg; r.font.name = FONT
    if sub:
        for line_txt in sub.split("\n"):
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = line_txt
            r2.font.size = Pt(ssize); r2.font.color.rgb = fg; r2.font.name = FONT
    return b


def arrow(slide, x, y, w, h, direction="right", color=TEAL_DK):
    shp = {"right": MSO_SHAPE.RIGHT_ARROW, "left": MSO_SHAPE.LEFT_ARROW,
           "down": MSO_SHAPE.DOWN_ARROW, "up": MSO_SHAPE.UP_ARROW}[direction]
    a = slide.shapes.add_shape(shp, x, y, w, h)
    _set_fill(a, color); _no_line(a)
    return a


def code_block(slide, x, y, w, h, code, size=12.5):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_fill(bg, CODE_BG); _no_line(bg)
    try:
        bg.adjustments[0] = 0.04
    except Exception:
        pass
    tf = bg.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(12); tf.margin_right = Pt(10)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    for i, cl in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        r = p.add_run(); r.text = cl if cl else " "
        r.font.name = FONT_MONO; r.font.size = Pt(size); r.font.color.rgb = CODE_FG
    return bg


def table(slide, x, y, w, h, headers, rows, col_widths=None,
          header_fill=NAVY, header_fg=WHITE, fsize=12, hsize=12.5, zebra=True):
    nrows = len(rows) + 1
    ncols = len(headers)
    gt = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = gt.table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(w * cw / total))
    for c, htext in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(7); cell.margin_right = Pt(5)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = htext
        r.font.bold = True; r.font.size = Pt(hsize); r.font.color.rgb = header_fg; r.font.name = FONT
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            fill = LIGHT if (zebra and ri % 2 == 0) else WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(7); cell.margin_right = Pt(5)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fsize); r.font.color.rgb = INK; r.font.name = FONT
    return tbl


def chip(slide, x, y, w, text, color, h=Inches(0.42), fg=WHITE, size=11.5):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_fill(c, color); _no_line(c)
    try:
        c.adjustments[0] = 0.5
    except Exception:
        pass
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = fg; r.font.name = FONT
    return c


def glyph_badge(slide, x, y, d, glyph, fill, size=20, fg=WHITE):
    """A round icon badge with a single glyph centred inside."""
    b = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    _set_fill(b, fill); _no_line(b)
    tf = b.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = glyph
    r.font.size = Pt(size); r.font.color.rgb = fg; r.font.name = FONT
    return b


def footer(slide, *_):
    n = len(prs.slides._sldIdLst)
    textbox(slide, Inches(11.4), Inches(7.06), Inches(1.7), Inches(0.3),
            [{"text": f"PingIt SDK   ·   {n}", "size": 8.5, "color": MUTED}],
            align=PP_ALIGN.RIGHT)


# ======================================================================
# SECTION 3 — SLIDES   (edit wording here)
# ======================================================================

def slide_title():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, NAVY); _no_line(bg)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.0), SLIDE_W, Inches(0.09))
    _set_fill(band, TEAL); _no_line(band)
    textbox(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.5),
            [{"text": "MOBILE SDK  ·  PROPOSAL & ARCHITECTURE", "size": 15, "bold": True, "color": TEAL}])
    textbox(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(1.5),
            [{"text": "PingIt SDK", "size": 56, "bold": True, "color": WHITE}])
    textbox(s, Inches(0.9), Inches(3.3), Inches(11.2), Inches(1.4),
            [{"text": "A lightweight SDK that measures real connection quality on mobile —",
              "size": 21, "color": PALE, "italic": True, "space_after": 2},
             {"text": "efficiently, and without depending on any single source.",
              "size": 21, "color": PALE, "italic": True}])
    textbox(s, Inches(0.9), Inches(5.45), Inches(11.0), Inches(1.1),
            [{"text": "Tom Bitran", "size": 18, "bold": True, "color": WHITE, "space_after": 3},
             {"text": "Mobile Development / SDK Class", "size": 14, "color": PALE}])
    # "ping" pulse motif — an on-brand visual (radar/ping rings)
    cxf, cyf = 11.5, 2.2
    for d, rc, wd in [(2.3, RGBColor(0x9F, 0xD9, 0xE8), 2.0),
                      (1.6, RGBColor(0x5F, 0xBF, 0xD6), 2.5),
                      (0.95, TEAL, 3.0)]:
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cxf - d / 2), Inches(cyf - d / 2), Inches(d), Inches(d))
        o.fill.background(); o.line.color.rgb = rc; o.line.width = Pt(wd)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cxf - 0.15), Inches(cyf - 0.15), Inches(0.3), Inches(0.3))
    _set_fill(dot, TEAL); _no_line(dot)
    return s


def slide_agenda():
    s = slide_base("What we'll cover", kicker="Agenda")
    items = [
        ("01", "The need & the idea"),
        ("02", "Features & readiness"),
        ("03", "System architecture"),
        ("04", "SDK: functions & data strategy"),
        ("05", "Server: functions & data model"),
        ("06", "Saving & retrieving efficiently"),
        ("07", "Request flow & resilience"),
        ("08", "Lightweight design & solutions"),
        ("09", "The portal"),
        ("10", "Research, stack & applications"),
    ]
    x0 = [Inches(0.7), Inches(6.85)]
    y0 = Inches(1.85)
    for i, (num, label) in enumerate(items):
        col = 0 if i < 5 else 1
        row = i % 5
        x = x0[col]; y = y0 + row * Inches(1.0)
        box(s, x, y, Inches(0.72), Inches(0.72), num, fill=TEAL, tsize=18)
        textbox(s, x + Inches(0.95), y, Inches(4.9), Inches(0.72),
                [{"text": label, "size": 16, "bold": True, "color": NAVY}],
                anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_idea():
    s = slide_base("The idea", kicker="Concept", accent=TEAL)
    textbox(s, Inches(0.7), Inches(1.65), Inches(12.0), Inches(0.95),
            [{"text": "PingIt is a lightweight mobile SDK that measures the real quality of a device's "
                      "internet connection — download, upload, latency, jitter (and packet loss).",
              "size": 17, "color": INK, "line_spacing": 1.1}])
    cards = [
        ("Lightweight", TEAL, "Small footprint, short tests, battery- & data-aware. A few lines to integrate."),
        ("Smart data strategy", NAVY2, "Decide on-device. Cache thresholds. Store history only if the developer wants it."),
        ("Resilient by design", GREEN, "Our server + public fallback, cached/bundled profiles, and honest offline handling."),
    ]
    x = Inches(0.7); w = Inches(3.92); gap = Inches(0.14)
    for i, (t, col, body) in enumerate(cards):
        px = x + i * (w + gap)
        panel(s, px, Inches(2.75), w, Inches(2.05), fill=LIGHT, line=col, line_w=1.5)
        head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(2.75), w, Inches(0.55))
        _set_fill(head, col); _no_line(head)
        tf = head.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = t; rr.font.bold = True; rr.font.size = Pt(15); rr.font.color.rgb = WHITE
        textbox(s, px + Inches(0.22), Inches(3.45), w - Inches(0.42), Inches(1.25),
                [{"text": body, "size": 12.5, "color": INK, "line_spacing": 1.08}])
    # readiness as one capability (de-emphasised)
    rp = panel(s, Inches(0.7), Inches(5.1), Inches(12.0), Inches(1.55), fill=NAVY)
    textbox(s, Inches(0.95), Inches(5.22), Inches(6.0), Inches(1.35),
            [{"text": "Plus — optional readiness answers", "size": 15, "bold": True, "color": TEAL, "space_after": 5},
             {"text": "Ask about a use case and get pass/fail, so developers needn't interpret raw numbers. "
                      "One handy capability — not the whole product.", "size": 12.5, "color": PALE, "line_spacing": 1.08}])
    code_block(s, Inches(7.1), Inches(5.3), Inches(5.4), Inches(1.15),
               'pingIt.isReadyFor("VIDEO_CALL")\n        ->  pass / fail', size=14)
    footer(s)
    return s


def slide_problem():
    s = slide_base("The problem & the need", kicker="Why", accent=AMBER)
    textbox(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(0.55),
            [{"text": "Mobile networks are unpredictable — and apps can't tell what they can handle.",
              "size": 18, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}])
    cards = [
        ("1", AMBER, "Networks vary wildly", "cellular to wifi, moving between cells, weak signal, dead zones"),
        ("2", RED, "Apps are flying blind", "no lightweight \"can the network handle this?\" check exists"),
        ("3", NAVY2, "Hard to do well", "slow-start, flaky nets, data cost, and you need a target to measure against"),
    ]
    wf, gapf, x0f = 3.92, 0.14, 0.7
    for i, (glyph, col, title, body) in enumerate(cards):
        pxf = x0f + i * (wf + gapf)
        panel(s, Inches(pxf), Inches(2.5), Inches(wf), Inches(3.0), fill=LIGHT, line=col, line_w=1.5)
        glyph_badge(s, Inches(pxf + (wf - 0.95) / 2), Inches(2.85), Inches(0.95), glyph, col, size=34)
        textbox(s, Inches(pxf + 0.2), Inches(4.05), Inches(wf - 0.4), Inches(0.5),
                [{"text": title, "size": 17, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}])
        textbox(s, Inches(pxf + 0.35), Inches(4.6), Inches(wf - 0.7), Inches(0.9),
                [{"text": body, "size": 13, "color": MUTED, "align": PP_ALIGN.CENTER, "line_spacing": 1.08}])
    strip = panel(s, Inches(0.7), Inches(5.78), Inches(12.0), Inches(0.82), fill=NAVY)
    textbox(s, Inches(0.95), Inches(5.8), Inches(11.6), Inches(0.78),
            [{"text": "PingIt gives them a lightweight read of the connection — and a clear answer — in a few lines.",
              "size": 14, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_features():
    s = slide_base("Features & capabilities", kicker="What it does", accent=TEAL)
    rows = [
        ["F1", "Speed test", "Download & upload throughput (Mbps)"],
        ["F2", "Latency test", "Round-trip time (ms) via repeated pings"],
        ["F3", "Jitter", "Variation between consecutive pings (RFC 3550)"],
        ["F4", "Packet loss", "% of probe packets lost (UDP/WebRTC enhancement)"],
        ["F5", "Readiness checks", "Pass/fail against named use-case profiles"],
        ["F6", "Quality score", "Optional 0–100 score + label for raw output"],
        ["F7", "Configurable history", "Dev chooses: none / last-in-cache / full server history"],
        ["F8", "Resilient sources", "Public fallback + explicit offline handling"],
    ]
    table(s, Inches(0.7), Inches(1.8), Inches(12.0), Inches(4.5),
          ["#", "Feature", "Description"], rows,
          col_widths=[1.1, 3.4, 7.8], fsize=14, hsize=14)
    textbox(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
            [{"text": "Lightweight & efficient by design — short tests, capped data, battery-aware, optional data-saver mode.",
              "size": 12.5, "italic": True, "color": MUTED}])
    footer(s)
    return s


def slide_readiness():
    s = slide_base("Readiness checks", kicker="One capability", accent=GREEN)
    textbox(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(0.55),
            [{"text": "Ask about a use case; PingIt compares the measurement to that profile's thresholds and returns pass/fail.",
              "size": 14, "color": INK}])
    rows = [
        ["MESSAGING", "Any working link"],
        ["WEB_BROWSING", "Modest download, low latency"],
        ["VOICE_CALL", "Low latency & jitter"],
        ["VIDEO_CALL", "Higher up+down, low latency"],
        ["HD_STREAMING", "Higher download"],
        ["UHD_4K_STREAMING", "High sustained download"],
        ["CLOUD_GAMING", "Ultra-low latency & jitter"],
        ["LIVE_BROADCAST", "High upload + low latency"],
    ]
    table(s, Inches(0.7), Inches(2.2), Inches(6.3), Inches(4.3),
          ["Profile", "Roughly needs"], rows,
          col_widths=[3.0, 3.3], fsize=12, hsize=12.5)
    code = (
        "val ready = pingIt.isReadyFor(\n"
        "    Profile.CLOUD_GAMING)\n\n"
        "if (!ready.passed)\n"
        "    warn(ready.reason)\n"
        "    // \"latency too high\"\n"
        "    //  or \"no connection\""
    )
    code_block(s, Inches(7.25), Inches(2.2), Inches(5.45), Inches(2.1), code, size=13)
    textbox(s, Inches(7.25), Inches(4.45), Inches(5.45), Inches(0.35),
            [{"text": "A check returns:", "size": 13, "bold": True, "color": NAVY}])
    ret = ("{ profile, passed, reason,\n"
           "  measured: { download, upload,\n"
           "    latency, jitter, packetLoss } }")
    code_block(s, Inches(7.25), Inches(4.85), Inches(5.45), Inches(1.5), ret, size=12.5)
    textbox(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
            [{"text": "Thresholds live server-side (tunable without an app release); the SDK caches them.",
              "size": 12, "italic": True, "color": MUTED}])
    footer(s)
    return s


def slide_architecture():
    s = slide_base("System architecture", kicker="The big picture", accent=NAVY2)
    textbox(s, Inches(0.7), Inches(1.58), Inches(12.0), Inches(0.5),
            [{"text": "Read left → right:  the phone measures & decides;  our server is just the target + a tiny config/data API.",
              "size": 13.5, "color": INK}])
    # phone
    box(s, Inches(0.55), Inches(2.35), Inches(2.85), Inches(2.45), "On the phone",
        "PingIt SDK\n\nmeasures the link\n+ decides pass/fail", fill=NAVY, tsize=15, ssize=11.5)
    textbox(s, Inches(0.55), Inches(4.92), Inches(2.85), Inches(0.7),
            [{"text": "All logic runs here — works on weak links.", "size": 10.5, "italic": True,
              "color": MUTED, "align": PP_ALIGN.CENTER}])
    arrow(s, Inches(3.45), Inches(3.3), Inches(0.6), Inches(0.55), "right")
    # server panel
    panel(s, Inches(4.15), Inches(2.05), Inches(5.55), Inches(3.3), fill=LIGHT, line=NAVY2, line_w=1.5)
    textbox(s, Inches(4.3), Inches(2.13), Inches(5.25), Inches(0.4),
            [{"text": "Measurement Server (ours)", "size": 13.5, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}])
    box(s, Inches(4.35), Inches(2.62), Inches(2.5), Inches(2.55), "Measurement",
        "/download\n/upload\n/ping\n/probe\n\nheavy bytes\nNO database", fill=TEAL_DK, tsize=13, ssize=10.5)
    box(s, Inches(7.05), Inches(2.62), Inches(2.5), Inches(2.55), "API / data",
        "/profiles\n/results\n/health\n\n\ntiny JSON\ntalks to DB", fill=NAVY2, tsize=13, ssize=10.5)
    arrow(s, Inches(9.75), Inches(3.3), Inches(0.55), Inches(0.55), "right")
    # postgres
    box(s, Inches(10.4), Inches(2.45), Inches(2.5), Inches(2.1), "Postgres",
        "profiles (always)\n\nresults (optional)", fill=NAVY, tsize=15, ssize=11.5)
    # bottom strip
    strip = panel(s, Inches(0.55), Inches(5.75), Inches(12.35), Inches(0.95), fill=LIGHT2)
    textbox(s, Inches(0.8), Inches(5.83), Inches(11.9), Inches(0.85),
            [{"text": "Public fallback (Cloudflare / M-Lab NDT) — used only if our server is down.", "size": 12.5, "color": NAVY, "bullet": True, "space_after": 2},
             {"text": "Web portal (login-gated) — edits profile thresholds and reads analytics.", "size": 12.5, "color": NAVY, "bullet": True}])
    footer(s)
    return s


def slide_sdk_exposed():
    s = slide_base("SDK — exposed (developer-facing) functions", kicker="Library API", accent=TEAL)
    textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.4),
            [{"text": "The small, stable surface a developer calls. Naming finalised per platform in design.",
              "size": 13.5, "color": MUTED}])
    code = (
        "// ---- Setup (appId scopes all data to this app) ------------------\n"
        "PingIt.init(appId, { historyMode, dataSaver, endpoint })\n\n"
        "// ---- Readiness --------------------------------------------------\n"
        "pingIt.isReadyFor(Profile.VIDEO_CALL)  -> { passed, reason, measured }\n\n"
        "// ---- Raw test ---------------------------------------------------\n"
        "pingIt.runTest()   -> { downloadMbps, uploadMbps, latencyMs, jitterMs,\n"
        "                        packetLossPct, score, label, timestamp }\n"
        "pingIt.ping()      -> latencyMs\n\n"
        "// ---- History (only when historyMode != NONE) --------------------\n"
        "pingIt.getHistory(limit)   -> [ past results ]\n\n"
        "// ---- Maintenance ------------------------------------------------\n"
        "pingIt.refreshProfiles()   // force a threshold refresh\n"
        "pingIt.cancel()            // abort an in-flight test"
    )
    code_block(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.55), code, size=13)
    footer(s)
    return s


def slide_sdk_internal():
    s = slide_base("SDK — internal functions", kicker="Under the hood", accent=NAVY2)
    textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.4),
            [{"text": "Hidden from the developer; this is where the measuring & decision logic lives.",
              "size": 13.5, "color": MUTED}])
    cols = [
        ("Measurement engine", TEAL_DK, [
            "measureDownload() — stream & time bytes",
            "measureUpload() — push & time bytes",
            "measurePing(n) — RTT samples",
            "computeJitter(samples) — RFC 3550",
            "warmupSkip() — drop TCP slow-start",
            "detectOffline() — no link? stop early",
        ]),
        ("Decision & data", NAVY2, [
            "evaluate(profile, measured) — pass/fail",
            "computeScore(measured) — 0–100",
            "selectTarget() — ours → public fallback",
            "adaptivePayload() — size by link speed",
        ]),
        ("Caching & transport", GREEN, [
            "profileCache — TTL + version + jitter",
            "loadBundledDefaults() — offline profiles",
            "lastResultCache — for LAST_LOCAL",
            "outbox — buffer + retry result upload",
            "httpClient — retries, timeouts",
        ]),
    ]
    x = Inches(0.7); w = Inches(3.92); gap = Inches(0.14)
    for i, (title, col, items) in enumerate(cols):
        px = x + i * (w + gap)
        titled_panel(s, px, Inches(2.15), w, Inches(4.4), title, col,
                     [(it, True) for it in items], body_size=11)
    footer(s)
    return s


def slide_server_functions():
    s = slide_base("Server — functions & endpoints", kicker="What the server does", accent=TEAL_DK)
    textbox(s, Inches(0.7), Inches(1.58), Inches(12.0), Inches(0.45),
            [{"text": "The server is deliberately “dumb”: it serves bytes and stores/returns JSON. "
                      "No decision logic lives here — that's on the phone.", "size": 13.5, "color": INK}])
    m = panel(s, Inches(0.7), Inches(2.15), Inches(5.9), Inches(4.4), fill=LIGHT, line=TEAL_DK, line_w=1.5)
    textbox(s, Inches(0.95), Inches(2.28), Inches(5.4), Inches(0.6),
            [{"text": "Measurement endpoints", "size": 16, "bold": True, "color": TEAL_DK, "space_after": 2},
             {"text": "high-bandwidth · short bursts · NO database", "size": 11, "color": MUTED}])
    textbox(s, Inches(0.95), Inches(3.05), Inches(5.45), Inches(3.4),
            [{"text": "GET /download?bytes=N", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "streams N incompressible bytes → download Mbps (size capped)", "size": 11.5, "color": MUTED, "level": 1, "space_after": 9},
             {"text": "POST /upload", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "reads & discards N bytes (sink) → upload Mbps", "size": 11.5, "color": MUTED, "level": 1, "space_after": 9},
             {"text": "GET /ping", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "tiny 200; called ×10 → latency + jitter", "size": 11.5, "color": MUTED, "level": 1, "space_after": 9},
             {"text": "/probe   (enhancement)", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "UDP/WebRTC echo → packet loss %", "size": 11.5, "color": MUTED, "level": 1}])
    a = panel(s, Inches(6.8), Inches(2.15), Inches(5.9), Inches(4.4), fill=LIGHT, line=NAVY2, line_w=1.5)
    textbox(s, Inches(7.05), Inches(2.28), Inches(5.4), Inches(0.6),
            [{"text": "API / data endpoints", "size": 16, "bold": True, "color": NAVY2, "space_after": 2},
             {"text": "low-bandwidth JSON · talks to Postgres", "size": 11, "color": MUTED}])
    textbox(s, Inches(7.05), Inches(3.05), Inches(5.45), Inches(3.4),
            [{"text": "GET /profiles", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "all profiles + version; SDK caches & refetches on change", "size": 11.5, "color": MUTED, "level": 1, "space_after": 9},
             {"text": "POST /results", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "validates appId, stamps time, stores (SERVER_HISTORY only)", "size": 11.5, "color": MUTED, "level": 1, "space_after": 9},
             {"text": "GET /results?appId&deviceId&limit", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "last N results for that device, scoped by app", "size": 11.5, "color": MUTED, "level": 1, "space_after": 9},
             {"text": "GET /health", "size": 13.5, "bold": True, "font": FONT_MONO, "color": NAVY, "space_after": 1},
             {"text": "200 OK — fallback logic + monitoring", "size": 11.5, "color": MUTED, "level": 1}])
    footer(s)
    return s


def slide_server_objects():
    s = slide_base("Server — object structure & data model", kicker="Postgres schema", accent=NAVY)
    textbox(s, Inches(0.7), Inches(1.62), Inches(6), Inches(0.35),
            [{"text": "profiles — tunable thresholds", "size": 14, "bold": True, "color": NAVY}])
    table(s, Inches(0.7), Inches(2.0), Inches(5.9), Inches(1.85),
          ["Column", "Type", "Notes"],
          [["id", "text (PK)", "e.g. VIDEO_CALL"],
           ["requires", "jsonb", "min/max per metric"],
           ["version", "int", "bumped on change"],
           ["updated_at", "timestamptz", ""]],
          col_widths=[2.0, 2.0, 2.4], fsize=11, hsize=11.5)
    textbox(s, Inches(0.7), Inches(4.0), Inches(6), Inches(0.35),
            [{"text": "apps — registered developer apps", "size": 14, "bold": True, "color": NAVY}])
    table(s, Inches(0.7), Inches(4.38), Inches(5.9), Inches(1.35),
          ["Column", "Type", "Notes"],
          [["app_id", "text (PK)", "issued in portal"],
           ["name", "text", "display name"],
           ["created_at", "timestamptz", ""]],
          col_widths=[2.0, 2.0, 2.4], fsize=11, hsize=11.5)
    textbox(s, Inches(6.85), Inches(1.62), Inches(6), Inches(0.35),
            [{"text": "results — history (SERVER_HISTORY only)", "size": 14, "bold": True, "color": NAVY}])
    table(s, Inches(6.85), Inches(2.0), Inches(5.85), Inches(3.7),
          ["Column", "Type", "Notes"],
          [["id", "bigserial PK", ""],
           ["app_id", "text", "scopes to an app"],
           ["device_id", "text", "anonymous, persists"],
           ["download/upload_mbps", "numeric", ""],
           ["latency/jitter_ms", "numeric", ""],
           ["packet_loss_pct", "numeric", ""],
           ["score", "int", "optional"],
           ["created_at", "timestamptz", "indexed"]],
          col_widths=[2.6, 1.7, 2.0], fsize=10.5, hsize=11)
    strip = panel(s, Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.7), fill=LIGHT2)
    textbox(s, Inches(0.9), Inches(5.99), Inches(11.7), Inches(0.62),
            [{"text": "Index  results(app_id, device_id, created_at DESC)  → “last N for this device” with no sort/scan.   "
                      "jsonb “requires” adds new metrics with no migration.",
              "size": 11.5, "color": NAVY}], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_data_efficiency():
    s = slide_base("Saving & retrieving data efficiently — and why", kicker="Server-side architecture", accent=GREEN)
    textbox(s, Inches(0.7), Inches(1.58), Inches(12.0), Inches(0.45),
            [{"text": "Each path is tuned for its one job: store cheaply, and read the common case for almost free.",
              "size": 13.5, "color": INK}])
    save_items = [
        ("Measurement traffic never touches the DB",
         "heavy MB bursts would swamp Postgres — streaming is network-bound, kept separate"),
        ("results is append-only (INSERT only)",
         "no UPDATE/DELETE → no row locks or contention; the index stays healthy"),
        ("History is opt-in (LAST_LOCAL by default)",
         "most apps write nothing server-side; storage grows only when trends are wanted"),
        ("Writes batched from the SDK outbox",
         "a closing/offline app never blocks; fewer, larger requests"),
    ]
    retrieve_items = [
        ("Composite index (app_id, device_id, created_at DESC)",
         "the only query is “last N, newest-first” — served directly: no scan, no sort"),
        ("limit-based pagination",
         "bounded, predictable reads — never “fetch everything”"),
        ("Profiles cached on-device (version-gated, jittered daily)",
         "thresholds are read constantly but change rarely → we almost never query the DB"),
        ("jsonb threshold object",
         "whole profile in one row; add a new metric with no schema migration"),
    ]

    def build_panel(px, title, color, items):
        panel(s, px, Inches(2.1), Inches(5.9), Inches(3.75), fill=LIGHT, line=color, line_w=1.5)
        head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(2.1), Inches(5.9), Inches(0.55))
        _set_fill(head, color); _no_line(head)
        tf = head.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = title
        rr.font.bold = True; rr.font.size = Pt(14); rr.font.color.rgb = WHITE; rr.font.name = FONT
        lines = []
        for dec, why in items:
            lines.append({"text": dec, "size": 12, "bold": True, "color": NAVY, "bullet": True, "space_after": 1})
            lines.append({"text": why, "size": 10.5, "color": MUTED, "level": 1, "space_after": 8, "line_spacing": 1.0})
        textbox(s, px + Inches(0.22), Inches(2.78), Inches(5.46), Inches(3.0), lines)

    build_panel(Inches(0.7), "Saving data  (write path)", TEAL_DK, save_items)
    build_panel(Inches(6.8), "Retrieving data  (read path)", GREEN, retrieve_items)
    strip = panel(s, Inches(0.7), Inches(5.98), Inches(12.0), Inches(0.65), fill=NAVY)
    textbox(s, Inches(0.95), Inches(6.0), Inches(11.6), Inches(0.6),
            [{"text": "Principle:  read the common case for free (cache),  write only when asked (opt-in),  index for the one query we actually run.",
              "size": 13, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_efficiency():
    s = slide_base("On-device data strategy", kicker="Caching · history · identity", accent=GREEN)
    textbox(s, Inches(0.7), Inches(1.58), Inches(12.0), Inches(0.45),
            [{"text": "Light by default: decide on-device, cache thresholds, and store history only if the developer asks.",
              "size": 13.5, "color": INK}])
    titled_panel(s, Inches(0.7), Inches(2.15), Inches(5.9), Inches(2.0),
                 "Cached on the device", TEAL_DK, [
                     "Profiles (+ their version)",
                     "The last result (for LAST_LOCAL)",
                     "A bounded outbox of unsent results",
                 ], body_size=12.5)
    titled_panel(s, Inches(6.8), Inches(2.15), Inches(5.9), Inches(2.0),
                 "Identity without accounts", NAVY2, [
                     "appId scopes all data to one app",
                     "deviceId = UUID in private storage; survives restarts",
                     "IP changing doesn't matter (we key on appId+deviceId)",
                 ], body_size=12.5)
    # history modes mini-table inside a titled panel area
    titled_panel(s, Inches(0.7), Inches(4.35), Inches(5.9), Inches(2.2),
                 "History mode — the developer chooses", AMBER, [], body_size=12)
    modes = [
        ("NONE", "keep nothing", RED),
        ("LAST_LOCAL  (default)", "last result in device cache only", GREEN),
        ("SERVER_HISTORY", "send every result to our server", NAVY2),
    ]
    for i, (mode, desc, col) in enumerate(modes):
        my = Inches(5.0) + i * Inches(0.5)
        chip(s, Inches(0.95), my, Inches(2.3), mode, col, h=Inches(0.4), size=10.5)
        textbox(s, Inches(3.35), my, Inches(3.1), Inches(0.4),
                [{"text": desc, "size": 11, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
    titled_panel(s, Inches(6.8), Inches(4.35), Inches(5.9), Inches(2.2),
                 "Refresh & server-down", GREEN, [
                     "Profiles refresh ~daily at a RANDOM time (no stampede)",
                     "Unsent results buffered + retried in the background",
                     "Server down? readiness still works from cached profiles",
                 ], body_size=12.5)
    footer(s)
    return s


def slide_flow():
    s = slide_base("Request flow — a CLOUD_GAMING check", kicker="End to end", accent=TEAL)
    steps = [
        ("1", "Fetch profiles", "GET /profiles\ncached (~daily, random)", TEAL_DK),
        ("2", "App asks", "isReadyFor(\nCLOUD_GAMING)", NAVY2),
        ("3", "Measure", "/download, /upload,\n/ping ×10", TEAL_DK),
        ("4", "Compare on-device", "measured vs cached\nthresholds", GREEN),
        ("5", "Return", "{ passed:false,\nreason:\"latency...\" }", NAVY),
        ("6", "Store (optional)", "if SERVER_HISTORY →\nPOST /results (buffered)", NAVY2),
    ]
    x = Inches(0.55); w = Inches(1.95); gap = Inches(0.12); y = Inches(2.75)
    for i, (num, title, sub, col) in enumerate(steps):
        px = x + i * (w + gap)
        box(s, px, y, w, Inches(1.7), title, sub, fill=col, tsize=12.5, ssize=9.5)
        box(s, px + Inches(0.02), y - Inches(0.35), Inches(0.5), Inches(0.5), num, fill=AMBER, tsize=14)
        if i < len(steps) - 1:
            arrow(s, px + w - Inches(0.02), y + Inches(0.6), gap + Inches(0.18), Inches(0.45), "right", color=MUTED)
    off = panel(s, Inches(0.55), Inches(4.95), Inches(12.35), Inches(0.85), fill=LIGHT, line=RED, line_w=1.5)
    textbox(s, Inches(0.8), Inches(4.99), Inches(11.9), Inches(0.78),
            [{"text": "No connection at all  →  can't measure  →  return { passed:false, reason:\"no connection\" }.  Never a fake pass.",
              "size": 13.5, "bold": True, "color": RED}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
            [{"text": "Why decide on-device (step 4)? The connection under test may be weak — the answer must not depend on a server round-trip.",
              "size": 12.5, "italic": True, "color": MUTED}])
    footer(s)
    return s


def slide_resilience():
    s = slide_base("Resilience & no single source", kicker="Fallback logic", accent=GREEN)
    textbox(s, Inches(0.7), Inches(1.62), Inches(12), Inches(0.35),
            [{"text": "Measurement target selection", "size": 14, "bold": True, "color": NAVY}])
    box(s, Inches(0.7), Inches(2.05), Inches(3.2), Inches(0.95), "1 · Our server", "GET /health OK?", fill=NAVY, tsize=13, ssize=10.5)
    arrow(s, Inches(3.95), Inches(2.28), Inches(0.7), Inches(0.5), "right", color=MUTED)
    box(s, Inches(4.7), Inches(2.05), Inches(3.2), Inches(0.95), "2 · Public endpoint", "Cloudflare / NDT", fill=TEAL_DK, tsize=13, ssize=10.5)
    arrow(s, Inches(7.95), Inches(2.28), Inches(0.7), Inches(0.5), "right", color=MUTED)
    box(s, Inches(8.7), Inches(2.05), Inches(4.0), Inches(0.95), "Method-based result", "valid numbers either way", fill=GREEN, tsize=13, ssize=10.5)
    textbox(s, Inches(0.7), Inches(3.2), Inches(12), Inches(0.35),
            [{"text": "Profiles (thresholds) selection", "size": 14, "bold": True, "color": NAVY}])
    box(s, Inches(0.7), Inches(3.6), Inches(3.2), Inches(0.95), "1 · /profiles", "fresh from server", fill=NAVY, tsize=13, ssize=10.5)
    arrow(s, Inches(3.95), Inches(3.83), Inches(0.7), Inches(0.5), "right", color=MUTED)
    box(s, Inches(4.7), Inches(3.6), Inches(3.2), Inches(0.95), "2 · Cached", "last good copy", fill=TEAL_DK, tsize=13, ssize=10.5)
    arrow(s, Inches(7.95), Inches(3.83), Inches(0.7), Inches(0.5), "right", color=MUTED)
    box(s, Inches(8.7), Inches(3.6), Inches(4.0), Inches(0.95), "3 · Bundled defaults", "ship-with-SDK values", fill=GREEN, tsize=13, ssize=10.5)
    off = panel(s, Inches(0.7), Inches(4.85), Inches(12.0), Inches(0.9), fill=LIGHT, line=RED, line_w=1.5)
    textbox(s, Inches(0.95), Inches(4.9), Inches(11.6), Inches(0.82),
            [{"text": "No connection at all  →  there is nothing to measure (it isn't “0”, it's unknown).  "
                      "PingIt returns reason:\"no connection\" — it does NOT fake a pass.", "size": 13, "bold": True, "color": RED, "line_spacing": 1.05}],
            anchor=MSO_ANCHOR.MIDDLE)
    note = panel(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.75), fill=LIGHT2)
    textbox(s, Inches(0.95), Inches(5.94), Inches(11.6), Inches(0.67),
            [{"text": "So: with a live connection, checks keep working even if our server is down (cached/bundled profiles). "
                      "With no connection, we report offline honestly. No single point of failure.", "size": 12, "color": NAVY, "line_spacing": 1.05}],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_mobile():
    s = slide_base("Lightweight & mobile-aware design", kicker="Battery · data · flaky nets", accent=TEAL)
    items = [
        ("Adaptive, capped payloads", "Start /download small, grow only on fast links; cap total test data to a few MB.", "1", TEAL),
        ("Short tests", "A few seconds total: quick bursts, not background polling.", "2", AMBER),
        ("Warm-up awareness", "Ignore the first moments so TCP slow-start doesn't make fast links look slow.", "3", GREEN),
        ("Battery-aware", "Runs only when the developer asks; no idle drain.", "4", TEAL_DK),
        ("Cache + jittered refresh", "Profiles fetched ~daily at a random time, not per check; saves data and battery.", "5", NAVY2),
        ("Data-saver mode", "Optional smaller payloads / shorter tests for metered plans; can follow the OS flag.", "6", TEAL),
    ]
    x = [Inches(0.7), Inches(6.75)]
    y0 = Inches(1.9)
    for i, (h, body, glyph, gcol) in enumerate(items):
        col = 0 if i < 3 else 1
        row = i % 3
        px = x[col]; py = y0 + row * Inches(1.55)
        glyph_badge(s, px, py, Inches(0.58), glyph, gcol, size=21)
        textbox(s, px + Inches(0.78), py - Inches(0.02), Inches(5.05), Inches(1.5),
                [{"text": h, "size": 14.5, "bold": True, "color": NAVY, "space_after": 2},
                 {"text": body, "size": 11.5, "color": MUTED, "line_spacing": 1.05}])
    footer(s)
    return s


def slide_solutions():
    s = slide_base("Solutions — data cost & abuse", kicker="Keeping it cheap & safe", accent=AMBER)
    titled_panel(s, Inches(0.7), Inches(1.95), Inches(5.9), Inches(4.5),
                 "Mobile data cost", AMBER, [
                     "Adaptive, capped payloads — a few MB max per test",
                     "Data-saver mode: smaller payloads / shorter tests",
                     "Developer can wire it to the OS data-saver flag",
                     "Short tests — quick bursts, never background polling",
                 ], body_size=13.5)
    titled_panel(s, Inches(6.8), Inches(1.95), Inches(5.9), Inches(4.5),
                 "Abuse protection", RED, [
                     "API key (appId) required — unknown callers rejected",
                     "Rate limiting per appId and per IP",
                     "Payload caps — server refuses oversized N",
                     "Monitoring via /health to catch abnormal bandwidth",
                 ], body_size=13.5)
    footer(s)
    return s


def slide_portal_functions():
    s = slide_base("The portal — functions & access", kicker="Login-gated web app", accent=NAVY2)
    cols = [
        ("Profile management", TEAL_DK, [
            "View / edit thresholds",
            "Create & retire profiles",
            "Bump version + publish",
            "Diff & rollback",
        ]),
        ("Analytics & history", NAVY2, [
            "Trends: speed / latency over time",
            "Pass-rate per profile",
            "Per-app (scoped) views",
            "Export CSV / JSON",
        ]),
        ("Apps, health & ops", GREEN, [
            "Register apps → issue appId / API key",
            "Server / endpoint uptime",
            "Bandwidth & fallback rate",
            "Crash & error reports",
        ]),
    ]
    x = Inches(0.7); w = Inches(3.92); gap = Inches(0.14)
    for i, (title, col, items) in enumerate(cols):
        px = x + i * (w + gap)
        titled_panel(s, px, Inches(1.9), w, Inches(3.5), title, col,
                     items, body_size=12)
    acc = panel(s, Inches(0.7), Inches(5.55), Inches(12.0), Inches(1.1), fill=LIGHT2)
    textbox(s, Inches(0.95), Inches(5.62), Inches(11.6), Inches(1.0),
            [{"text": "Who can enter, and how", "size": 13, "bold": True, "color": NAVY, "space_after": 3},
             {"text": "PingIt operators / admins  sign in to edit profiles & watch health.   "
                      "Registered developers  sign in to get an appId and see THEIR app's analytics only.", "size": 12, "color": INK, "space_after": 1, "bullet": True},
             {"text": "Access is scoped by role and by appId — a developer never sees another app's data.", "size": 12, "color": INK, "bullet": True}])
    footer(s)
    return s


def slide_portal_wireframe():
    s = slide_base("The portal — schematic wireframe", kicker="Drawing", accent=TEAL)
    panel(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(4.95), fill=WHITE, line=NAVY2, line_w=1.5)
    topbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.75), Inches(12.0), Inches(0.5))
    _set_fill(topbar, NAVY); _no_line(topbar)
    textbox(s, Inches(0.95), Inches(1.75), Inches(8), Inches(0.5),
            [{"text": "PingIt · Admin Portal", "size": 12.5, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(9.5), Inches(1.75), Inches(3.0), Inches(0.5),
            [{"text": "signed in: admin", "size": 10.5, "color": PALE, "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.25), Inches(2.4), Inches(4.45))
    _set_fill(side, LIGHT2); _no_line(side)
    nav = ["Dashboard", "Profiles", "History", "Apps & keys", "Health", "Crashes"]
    for i, n in enumerate(nav):
        panel(s, Inches(0.85), Inches(2.45) + i * Inches(0.62), Inches(2.1), Inches(0.48),
              fill=(TEAL if i == 1 else WHITE), line=LIGHT2)
        textbox(s, Inches(1.0), Inches(2.45) + i * Inches(0.62), Inches(1.95), Inches(0.48),
                [{"text": n, "size": 11.5, "bold": (i == 1), "color": (WHITE if i == 1 else INK)}],
                anchor=MSO_ANCHOR.MIDDLE)
    kpis = [("Tests today", "12,480"), ("Pass rate", "78%"), ("Avg latency", "46 ms"), ("Fallback", "2.1%")]
    for i, (k, v) in enumerate(kpis):
        kx = Inches(3.35) + i * Inches(2.32)
        panel(s, kx, Inches(2.45), Inches(2.15), Inches(1.1), fill=LIGHT)
        textbox(s, kx + Inches(0.15), Inches(2.55), Inches(1.9), Inches(1.0),
                [{"text": v, "size": 22, "bold": True, "color": NAVY, "space_after": 0},
                 {"text": k, "size": 10.5, "color": MUTED}])
    panel(s, Inches(3.35), Inches(3.75), Inches(4.5), Inches(2.7), fill=LIGHT)
    textbox(s, Inches(3.55), Inches(3.85), Inches(4.1), Inches(0.4),
            [{"text": "Profile editor  ·  VIDEO_CALL", "size": 12, "bold": True, "color": NAVY}])
    rows_lbl = ["downloadMbps  min  1.5", "uploadMbps    min  1.0", "latencyMs     max  200", "jitterMs      max  40"]
    for i, rl in enumerate(rows_lbl):
        panel(s, Inches(3.55), Inches(4.35) + i * Inches(0.46), Inches(4.1), Inches(0.38), fill=WHITE, line=LIGHT2)
        textbox(s, Inches(3.7), Inches(4.35) + i * Inches(0.46), Inches(3.9), Inches(0.38),
                [{"text": rl, "size": 10.5, "font": FONT_MONO, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
    panel(s, Inches(8.05), Inches(3.75), Inches(4.5), Inches(2.7), fill=LIGHT)
    textbox(s, Inches(8.25), Inches(3.85), Inches(4.1), Inches(0.4),
            [{"text": "Latency trend (7 days)", "size": 12, "bold": True, "color": NAVY}])
    heights = [0.7, 1.0, 0.6, 1.3, 0.9, 1.5, 1.1]
    for i, h in enumerate(heights):
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.35) + i * Inches(0.58),
                                 Inches(6.2) - Inches(h), Inches(0.42), Inches(h))
        _set_fill(bar, TEAL); _no_line(bar)
    footer(s)
    return s


def slide_research_android():
    s = slide_base("Research — data we can read from an Android device", kicker="Required research · device analytics", accent=NAVY2)
    textbox(s, Inches(0.7), Inches(1.58), Inches(12), Inches(0.45),
            [{"text": "What the platform lets us collect to understand network conditions (privacy-safe, no PII).",
              "size": 13.5, "color": INK}])
    cols = [
        ("Connectivity & network", TEAL_DK, [
            "ConnectivityManager /",
            "NetworkCapabilities:",
            "  WIFI / CELLULAR, metered flag,",
            "  up/down bandwidth hints",
            "TelephonyManager: LTE/5G,",
            "  carrier, SIM country",
        ]),
        ("Signal & radio", NAVY2, [
            "SignalStrength /",
            "CellSignalStrength:",
            "  RSRP, RSRQ, dBm, level 0-4",
            "WifiManager / WifiInfo:",
            "  RSSI, link speed,",
            "  frequency (2.4 / 5 GHz)",
        ]),
        ("Device & context", GREEN, [
            "Build: model, manufacturer,",
            "  API level",
            "Battery state, power-save mode",
            "App foreground / background",
            "Network-change callbacks",
        ]),
    ]
    x = Inches(0.7); w = Inches(3.92); gap = Inches(0.14)
    for i, (title, col, items) in enumerate(cols):
        px = x + i * (w + gap)
        panel(s, px, Inches(2.15), w, Inches(3.55), fill=LIGHT, line=col, line_w=1.5)
        head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(2.15), w, Inches(0.55))
        _set_fill(head, col); _no_line(head)
        tf = head.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = title; rr.font.bold = True; rr.font.size = Pt(12.5); rr.font.color.rgb = WHITE
        lines = [{"text": it, "size": 10.5, "color": INK, "font": FONT_MONO, "space_after": 4} for it in items]
        textbox(s, px + Inches(0.18), Inches(2.85), w - Inches(0.34), Inches(2.7), lines)
    note = panel(s, Inches(0.7), Inches(5.85), Inches(12.0), Inches(0.8), fill=LIGHT2)
    textbox(s, Inches(0.95), Inches(5.9), Inches(11.6), Inches(0.7),
            [{"text": "Privacy: no precise location, no PII. Anonymous deviceId only — analytics correlate network conditions with outcomes, not users.",
              "size": 12, "italic": True, "color": NAVY}], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_research_crash():
    s = slide_base("Research — catch bugs & upload before the app closes", kicker="Required research · reliability", accent=RED)
    textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.45),
            [{"text": "A crash or a swipe-away kills the process before data is sent. The fix is persist-first, upload-later.",
              "size": 13.5, "color": INK}])
    steps = [
        ("Catch", TEAL_DK, "Thread.setDefaultUncaught-\nExceptionHandler grabs\nuncaught exceptions + stack"),
        ("Persist now", NAVY2, "Write crash + last result to\nlocal disk SYNCHRONOUSLY\n(survives process death)"),
        ("Flush on stop", AMBER, "Lifecycle ON_STOP /\nonTrimMemory → best-effort\nflush of buffered results"),
        ("Retry upload", GREEN, "WorkManager re-sends on\nnext launch / network\n→ guaranteed delivery"),
    ]
    x = Inches(0.7); w = Inches(2.85); gap = Inches(0.2); y = Inches(2.45)
    for i, (title, col, body) in enumerate(steps):
        px = x + i * (w + gap)
        box(s, px, y, w, Inches(1.85), title, body, fill=col, tsize=15, ssize=10.5)
        if i < len(steps) - 1:
            arrow(s, px + w - Inches(0.02), y + Inches(0.65), gap + Inches(0.16), Inches(0.5), "right", color=MUTED)
    extra = [
        {"text": "Also handled:", "size": 14, "bold": True, "color": NAVY, "space_after": 6},
        {"text": "ANR detection (main-thread watchdog) and caught non-fatal errors logged as breadcrumbs", "size": 12.5, "color": INK, "bullet": True, "space_after": 5},
        {"text": "In-flight test state saved, so a test interrupted by a crash isn't reported as a false failure", "size": 12.5, "color": INK, "bullet": True, "space_after": 5},
        {"text": "Bounded local queue + exponential backoff so retries never drain battery or data", "size": 12.5, "color": INK, "bullet": True},
    ]
    textbox(s, Inches(0.7), Inches(4.7), Inches(12.0), Inches(2.0), extra)
    footer(s)
    return s


def slide_applications():
    s = slide_base("Applications that can use PingIt", kicker="Who benefits", accent=TEAL)
    apps = [
        ("Messaging / chat", "Check the link before sync; warn on dead links", TEAL_DK, "1"),
        ("VoIP / calling", "Warn before a weak voice or video call", NAVY2, "2"),
        ("Video streaming", "Pick HD vs SD from measured download", TEAL_DK, "3"),
        ("Cloud / file backup", "Defer big uploads until the link is strong", GREEN, "4"),
        ("Cloud gaming", "Gate on ultra-low latency and jitter", NAVY2, "5"),
        ("Live broadcasting", "Check upload before going live", TEAL_DK, "6"),
        ("Telehealth", "Video visits: check call quality first", GREEN, "7"),
        ("IoT / field apps", "Confirm the link before sending sensor data", NAVY2, "8"),
        ("Maps / navigation", "Weak data: switch to offline tiles", TEAL_DK, "9"),
    ]
    w = Inches(3.92); gap = Inches(0.14); h = Inches(1.4); badge_d = Inches(0.62)
    for i, (title, body, col, glyph) in enumerate(apps):
        c = i % 3; r = i // 3
        px = Inches(0.7) + c * (w + gap); py = Inches(1.85) + r * (h + Inches(0.18))
        panel(s, px, py, w, h, fill=LIGHT, line=col, line_w=1.25)
        glyph_badge(s, px + Inches(0.2), py + (h - badge_d) // 2, badge_d, glyph, col, size=20)
        textbox(s, px + Inches(1.0), py + Inches(0.16), w - Inches(1.2), h - Inches(0.3),
                [{"text": title, "size": 13.5, "bold": True, "color": NAVY, "space_after": 3},
                 {"text": body, "size": 11.5, "color": MUTED, "line_spacing": 1.05}])
    footer(s)
    return s


def slide_techstack():
    s = slide_base("Tech stack (proposed)", kicker="How we'll build it", accent=NAVY2)
    table(s, Inches(0.7), Inches(1.9), Inches(12.0), Inches(3.6),
          ["Concern", "Choice", "Why"],
          [["SDK", "iOS (Swift), Android (Kotlin)", "Shared core (Kotlin Multiplatform) TBD"],
           ["Measurement server", "Node.js", "Simple streaming; Go is the scale path"],
           ["API server", "Node.js", "Keep simple; can split out later"],
           ["Database", "Postgres", "jsonb thresholds + optional history"],
           ["Transport", "HTTPS (+ UDP/WebRTC)", "HTTPS for speed/ping; UDP for packet loss"],
           ["Hosting", "Cloud region near users", "Low, consistent measurement latency"]],
          col_widths=[2.6, 3.6, 5.8], fsize=13, hsize=13.5)
    strip = panel(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.95), fill=NAVY)
    textbox(s, Inches(0.95), Inches(5.78), Inches(11.6), Inches(0.85),
            [{"text": "Light, efficient, resilient — measure the connection, decide on-device, and store only what's needed.",
              "size": 14, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


# ======================================================================
# SECTION 4 — BUILD  (the running order of the deck)
# ======================================================================

def build():
    slide_title()
    slide_problem()          # open with the need...
    slide_idea()             # ...then the solution
    slide_features()
    slide_readiness()
    slide_architecture()
    slide_sdk_exposed()
    slide_sdk_internal()
    slide_efficiency()       # on-device data strategy — grouped with the SDK
    slide_server_functions()
    slide_server_objects()
    slide_data_efficiency()  # server save/retrieve — right after the data model
    slide_flow()
    slide_resilience()
    slide_mobile()
    slide_solutions()
    slide_portal_functions()
    slide_portal_wireframe()
    slide_research_android()
    slide_research_crash()
    slide_techstack()
    slide_applications()     # close on use cases

    import os
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "PingIt_SDK.pptx")
    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
