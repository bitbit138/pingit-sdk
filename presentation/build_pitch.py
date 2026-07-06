#!/usr/bin/env python3
"""
PingIt developer pitch deck (4-minute, non-technical, for the class-as-developers).

Separate from the detailed proposal deck (build_deck.py). Deliberately different
look: warm cream background, one bold orange accent, charcoal text, large
left-aligned numerals, thin rules, lots of negative space.

Regenerate:  ./.venv/bin/python presentation/build_pitch.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- palette ---
ORANGE = RGBColor(0xED, 0x6A, 0x2C)
ORANGE_DEEP = RGBColor(0xC2, 0x4E, 0x17)
INK = RGBColor(0x26, 0x1D, 0x17)      # warm near-black
CREAM = RGBColor(0xFB, 0xF3, 0xE9)    # background
CREAM2 = RGBColor(0xF2, 0xE6, 0xD5)   # subtle panel
MUTE = RGBColor(0x8A, 0x7C, 0x6E)     # warm gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color=CREAM):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(r, color)
    return r


def rect(slide, x, y, w, h, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(r, color)
    return r


def oval(slide, x, y, w, h, color):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    _fill(o, color)
    return o


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom):
        pass
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        if ln.get("line_spacing"):
            p.line_spacing = ln["line_spacing"]
        r = p.add_run()
        r.text = ln["text"]
        f = r.font
        f.size = Pt(ln.get("size", 18))
        f.bold = ln.get("bold", False)
        f.italic = ln.get("italic", False)
        f.name = ln.get("font", FONT)
        f.color.rgb = ln.get("color", INK)
        sp = ln.get("spacing")
        if sp is not None:
            _letter_spacing(r, sp)
    return tb


def _letter_spacing(run, pts):
    # set character spacing (tracking) in points via XML
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def kicker(slide, x, y, s, color=ORANGE_DEEP):
    text(slide, x, y, Inches(11), Inches(0.4),
         [{"text": s.upper(), "size": 13, "bold": True, "color": color, "spacing": 2.2}])


def pill(slide, x, y, w, h, label, fill, fg, size=14, outline=None):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is None:
        p.fill.background()
    else:
        p.fill.solid()
        p.fill.fore_color.rgb = fill
    if outline:
        p.line.color.rgb = outline
        p.line.width = Pt(1.5)
    else:
        p.line.fill.background()
    try:
        p.adjustments[0] = 0.5
    except Exception:
        pass
    tf = p.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = fg
    return p


def slide():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    return s


# ============================================================ 1. TITLE
s = slide()
oval(s, Inches(9.7), Inches(-1.6), Inches(5.6), Inches(5.6), CREAM2)      # soft motif
oval(s, Inches(11.2), Inches(0.2), Inches(2.6), Inches(2.6), ORANGE)      # bold dot
kicker(s, Inches(0.9), Inches(0.95), "Network-quality SDK  ·  for developers")
text(s, Inches(0.85), Inches(1.7), Inches(9.5), Inches(2.2),
     [{"text": "PingIt", "size": 96, "bold": True, "color": INK}])
rect(s, Inches(0.95), Inches(3.55), Inches(3.4), Inches(0.16), ORANGE)
text(s, Inches(0.9), Inches(3.95), Inches(9.6), Inches(1.6),
     [{"text": "Know the network before your app needs it.", "size": 30, "bold": True, "color": ORANGE_DEEP}])
text(s, Inches(0.9), Inches(6.5), Inches(9), Inches(0.5),
     [{"text": "A 4-minute intro for the developers who'd build with it.", "size": 15, "color": MUTE}])

# ============================================================ 2. PROBLEM
s = slide()
kicker(s, Inches(0.9), Inches(0.9), "The problem")
text(s, Inches(0.85), Inches(1.7), Inches(11.6), Inches(2.6),
     [{"text": "Your app can't tell if the connection can handle what the user is "
               "about to do.", "size": 40, "bold": True, "color": INK, "line_spacing": 1.05}])
text(s, Inches(0.9), Inches(4.5), Inches(11.4), Inches(1.6),
     [{"text": "So the video call freezes. The upload dies at 80%. The stream buffers. "
               "The live goes down in front of an audience.", "size": 22, "color": MUTE, "line_spacing": 1.2}])
rect(s, Inches(0.95), Inches(6.35), Inches(2.2), Inches(0.14), ORANGE)

# ============================================================ 3. THE VALUE LINE
s = slide()
kicker(s, Inches(0.9), Inches(0.9), "The idea, in one line")
text(s, Inches(0.85), Inches(1.55), Inches(11.6), Inches(0.9),
     [{"text": "You ask about the activity. PingIt answers yes or no.", "size": 30, "bold": True, "color": INK}])
# code chip
chip = rect(s, Inches(0.9), Inches(2.9), Inches(7.6), Inches(1.15), INK)
tf = chip.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Pt(18)
p = tf.paragraphs[0]
r = p.add_run()
r.text = 'pingIt.isReadyFor( VIDEO_CALL )'
r.font.name = MONO
r.font.size = Pt(24)
r.font.bold = True
r.font.color.rgb = CREAM
# arrow + outcomes
text(s, Inches(8.65), Inches(2.9), Inches(0.9), Inches(1.15),
     [{"text": "->", "size": 34, "bold": True, "color": ORANGE, "align": PP_ALIGN.CENTER}],
     anchor=MSO_ANCHOR.MIDDLE)
pill(s, Inches(9.5), Inches(3.02), Inches(3.3), Inches(0.44), "READY  -  go ahead", ORANGE, WHITE, size=13)
pill(s, Inches(9.5), Inches(3.55), Inches(3.6), Inches(0.44), "NOT READY  -  latency too high", None, INK, size=12, outline=INK)
text(s, Inches(0.9), Inches(4.9), Inches(11.4), Inches(1.6),
     [{"text": "No megabit math. No guessing what \"good enough\" means. One call, a clear "
               "answer, and the reason when it fails.", "size": 21, "color": MUTE, "line_spacing": 1.2}])

# ============================================================ 4. PROFILES
s = slide()
kicker(s, Inches(0.9), Inches(0.9), "The profiles feature")
text(s, Inches(0.85), Inches(1.55), Inches(11.6), Inches(1.4),
     [{"text": "Ask about what the user is doing, not about megabits.", "size": 32, "bold": True, "color": INK, "space_after": 6},
      {"text": "Pick from named use cases. We keep the thresholds tuned for you, server-side.",
       "size": 18, "color": MUTE}])
profiles = [
    ("MESSAGING", "any working link"),
    ("VIDEO CALL", "smooth two-way video"),
    ("HD & 4K STREAMING", "play without buffering"),
    ("CLOUD GAMING", "ultra-low latency"),
    ("GOING LIVE", "strong, steady upload"),
    ("BIG UPLOAD / BACKUP", "finish large transfers"),
]
y0 = 3.25
for i, (name, meaning) in enumerate(profiles):
    col = i // 3
    row = i % 3
    x = Inches(0.95 + col * 6.15)
    y = Inches(y0 + row * 1.18)
    rect(s, x, y + Inches(0.06), Inches(0.09), Inches(0.62), ORANGE)
    text(s, x + Inches(0.32), y, Inches(5.6), Inches(0.9),
         [{"text": name, "size": 19, "bold": True, "color": INK, "space_after": 1},
          {"text": meaning, "size": 14, "color": MUTE}])

# ============================================================ 5. THE DASHBOARD
s = slide()
kicker(s, Inches(0.9), Inches(0.7), "Behind it: the dashboard")
text(s, Inches(0.85), Inches(1.2), Inches(11.9), Inches(0.8),
     [{"text": "Tune the profiles live. No app update.", "size": 30, "bold": True, "color": INK}])
text(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(0.6),
     [{"text": "An operator manages thresholds, apps, and health in a web dashboard. "
               "Devices pick up the changes on their own.", "size": 16, "color": MUTE, "line_spacing": 1.15}])

# ---- dashboard mock (left) ----
fx, fy, fw, fh = 0.85, 2.75, 6.7, 4.35
frame = rect(s, Inches(fx), Inches(fy), Inches(fw), Inches(fh), WHITE)
frame.line.color.rgb = RGBColor(0xE2, 0xD6, 0xC6)
frame.line.width = Pt(1.25)
# window chrome
rect(s, Inches(fx), Inches(fy), Inches(fw), Inches(0.34), INK)
for j, c in enumerate((ORANGE, RGBColor(0xF2, 0xC0, 0x5A), RGBColor(0x9C, 0x8F, 0x80))):
    oval(s, Inches(fx + 0.18 + j * 0.22), Inches(fy + 0.1), Inches(0.13), Inches(0.13), c)
text(s, Inches(fx + 1.0), Inches(fy + 0.045), Inches(4), Inches(0.3),
     [{"text": "PingIt  ·  Dashboard", "size": 10, "color": CREAM}])
# sidebar
sbx, sby, sbw = fx, fy + 0.34, 1.75
sbh = fh - 0.34
rect(s, Inches(sbx), Inches(sby), Inches(sbw), Inches(sbh), CREAM2)
for k, item in enumerate(["Dashboard", "Profiles", "History", "Apps & keys", "Health", "Crashes"]):
    iy = sby + 0.3 + k * 0.62
    active = item == "Profiles"
    if active:
        rect(s, Inches(sbx), Inches(iy - 0.05), Inches(0.08), Inches(0.44), ORANGE)
        rect(s, Inches(sbx + 0.08), Inches(iy - 0.05), Inches(sbw - 0.08), Inches(0.44),
             RGBColor(0xEE, 0xDD, 0xC7))
    text(s, Inches(sbx + 0.26), Inches(iy), Inches(sbw - 0.34), Inches(0.4),
         [{"text": item, "size": 12.5, "bold": active,
           "color": (ORANGE_DEEP if active else INK)}], anchor=MSO_ANCHOR.MIDDLE)
# main panel label
mpx = fx + sbw
text(s, Inches(mpx + 0.3), Inches(sby + 0.06), Inches(fw - sbw - 0.5), Inches(0.3),
     [{"text": "PROFILES", "size": 10, "bold": True, "color": MUTE, "spacing": 1.8}])
# profile editor card
cx, cy, cw, ch = mpx + 0.28, sby + 0.42, fw - sbw - 0.55, sbh - 0.62
card = rect(s, Inches(cx), Inches(cy), Inches(cw), Inches(ch), CREAM)
card.line.color.rgb = RGBColor(0xE7, 0xD8, 0xC4)
card.line.width = Pt(1)
text(s, Inches(cx + 0.22), Inches(cy + 0.16), Inches(cw - 1.4), Inches(0.4),
     [{"text": "VIDEO CALL", "size": 17, "bold": True, "color": INK}])
pill(s, Inches(cx + cw - 0.95), Inches(cy + 0.18), Inches(0.7), Inches(0.34), "v7", ORANGE, WHITE, size=12)
rows = [("Download min", "1.5 Mbps"), ("Upload min", "1.0 Mbps"),
        ("Latency max", "200 ms"), ("Jitter max", "40 ms")]
ry0 = cy + 0.72
for r_i, (lbl, val) in enumerate(rows):
    ry = ry0 + r_i * 0.5
    text(s, Inches(cx + 0.24), Inches(ry), Inches(2.4), Inches(0.4),
         [{"text": lbl, "size": 13, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
    vb = rect(s, Inches(cx + cw - 1.55), Inches(ry - 0.01), Inches(1.3), Inches(0.4), WHITE)
    vb.line.color.rgb = RGBColor(0xD8, 0xC7, 0xB2)
    vb.line.width = Pt(1)
    vtf = vb.text_frame
    vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    vtf.margin_top = vtf.margin_bottom = Pt(0)
    vtf.margin_left = vtf.margin_right = Pt(6)
    vp = vtf.paragraphs[0]
    vp.alignment = PP_ALIGN.CENTER
    vr = vp.add_run()
    vr.text = val
    vr.font.size = Pt(13)
    vr.font.bold = True
    vr.font.name = MONO
    vr.font.color.rgb = INK
sbtn_y = cy + ch - 0.62
pill(s, Inches(cx + 0.24), Inches(sbtn_y), Inches(1.8), Inches(0.44), "Save changes", ORANGE, WHITE, size=13)
text(s, Inches(cx + 2.2), Inches(sbtn_y + 0.02), Inches(cw - 2.35), Inches(0.44),
     [{"text": "-> now version 8", "size": 12, "italic": True, "color": MUTE, "font": MONO}],
     anchor=MSO_ANCHOR.MIDDLE)

# ---- what you manage (right) ----
rx = 8.0
text(s, Inches(rx), Inches(2.78), Inches(4.8), Inches(0.5),
     [{"text": "What you manage here", "size": 19, "bold": True, "color": INK}])
manage = [
    "Tune profile thresholds - no app release needed.",
    "Register apps and issue an appId.",
    "Watch live analytics and device history.",
    "See crash reports as they arrive.",
]
by0 = 3.5
for b_i, btxt in enumerate(manage):
    byy = by0 + b_i * 0.6
    rect(s, Inches(rx), Inches(byy + 0.07), Inches(0.16), Inches(0.16), ORANGE)
    text(s, Inches(rx + 0.36), Inches(byy), Inches(4.5), Inches(0.6),
         [{"text": btxt, "size": 15, "color": INK, "line_spacing": 1.05}])
co_y = by0 + len(manage) * 0.6 + 0.15
rect(s, Inches(rx), Inches(co_y), Inches(0.1), Inches(0.9), ORANGE)
rect(s, Inches(rx + 0.1), Inches(co_y), Inches(4.7), Inches(0.9), CREAM2)
text(s, Inches(rx + 0.36), Inches(co_y + 0.14), Inches(4.3), Inches(0.7),
     [{"text": "Change what \"ready\" means from here - every device follows, no update.",
       "size": 13.5, "bold": True, "color": INK, "line_spacing": 1.1}])

# ============================================================ 6. HOW SIMPLE
s = slide()
kicker(s, Inches(0.9), Inches(0.9), "How you use it")
text(s, Inches(0.85), Inches(1.5), Inches(11.6), Inches(0.9),
     [{"text": "Three steps. That's the whole integration.", "size": 32, "bold": True, "color": INK}])
steps = [
    ("01", "Add the SDK", "It starts itself on app launch. No setup code."),
    ("02", "Ask", "pingIt.isReadyFor(Profile.VIDEO_CALL)"),
    ("03", "React", "Go ahead, or gently guide the user."),
]
for i, (num, title_, body) in enumerate(steps):
    y = Inches(2.85 + i * 1.35)
    text(s, Inches(0.95), y - Inches(0.15), Inches(1.5), Inches(1.1),
         [{"text": num, "size": 46, "bold": True, "color": ORANGE}])
    is_code = i == 1
    text(s, Inches(2.5), y, Inches(10), Inches(1.1),
         [{"text": title_, "size": 22, "bold": True, "color": INK, "space_after": 2},
          {"text": body, "size": 16, "color": MUTE, "font": (MONO if is_code else FONT)}])

def app_icon(slide_, x, y, d, brand):
    """A small app-icon-like tile: brand colour plus a simple mark."""
    def mk(shape_type, ox, oy, ow, oh):
        return slide_.shapes.add_shape(
            shape_type, x + int(d * ox), y + int(d * oy), int(d * ow), int(d * oh))

    tile = slide_.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, d, d)
    tile.line.fill.background()
    try:
        tile.adjustments[0] = 0.3
    except Exception:
        pass

    def fill(c):
        tile.fill.solid()
        tile.fill.fore_color.rgb = c

    def letter(txt, color, size):
        tf = tile.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT

    if brand == "netflix":
        fill(RGBColor(0x14, 0x14, 0x14))
        letter("N", RGBColor(0xE5, 0x09, 0x14), 34)
    elif brand == "xbox":
        fill(RGBColor(0x10, 0x7C, 0x10))
        letter("X", WHITE, 32)
    elif brand == "uber":
        fill(RGBColor(0x10, 0x10, 0x10))
        letter("Uber", WHITE, 15)
    elif brand == "whatsapp":
        fill(RGBColor(0x25, 0xD3, 0x66))
        b = mk(MSO_SHAPE.OVAL_CALLOUT, 0.18, 0.16, 0.62, 0.52)
        b.fill.solid()
        b.fill.fore_color.rgb = WHITE
        b.line.fill.background()
    elif brand == "instagram":
        fill(RGBColor(0xC1, 0x35, 0x84))
        body_ = mk(MSO_SHAPE.ROUNDED_RECTANGLE, 0.2, 0.28, 0.6, 0.44)
        body_.fill.background()
        body_.line.color.rgb = WHITE
        body_.line.width = Pt(2.2)
        try:
            body_.adjustments[0] = 0.3
        except Exception:
            pass
        lens = mk(MSO_SHAPE.OVAL, 0.37, 0.4, 0.26, 0.26)
        lens.fill.background()
        lens.line.color.rgb = WHITE
        lens.line.width = Pt(2.2)
        dot = mk(MSO_SHAPE.OVAL, 0.62, 0.32, 0.09, 0.09)
        dot.fill.solid()
        dot.fill.fore_color.rgb = WHITE
        dot.line.fill.background()
    elif brand == "googlephotos":
        fill(WHITE)
        tile.line.color.rgb = RGBColor(0xDE, 0xDE, 0xDE)
        tile.line.width = Pt(1)
        petals = [
            (0.34, 0.1, RGBColor(0xEA, 0x43, 0x35)),
            (0.56, 0.34, RGBColor(0xFB, 0xBC, 0x05)),
            (0.34, 0.56, RGBColor(0x34, 0xA8, 0x53)),
            (0.12, 0.34, RGBColor(0x42, 0x85, 0xF4)),
        ]
        for px, py, c in petals:
            pet = mk(MSO_SHAPE.OVAL, px, py, 0.32, 0.32)
            pet.fill.solid()
            pet.fill.fore_color.rgb = c
            pet.line.fill.background()
    else:
        fill(ORANGE)


# ============================================================ 7. SIX APPS
s = slide()
kicker(s, Inches(0.9), Inches(0.7), "Who would build with it")
text(s, Inches(0.85), Inches(1.25), Inches(11.6), Inches(0.7),
     [{"text": "Six apps, and how PingIt helps.", "size": 30, "bold": True, "color": INK}])
apps = [
    ("instagram", "Instagram", "Check before a big creator goes Live, so a shaky link never tanks the stream, or their reputation."),
    ("netflix", "Netflix", "Pick HD or SD up front, so playback starts instantly instead of buffering."),
    ("whatsapp", "WhatsApp", "Weak link before a video call? Offer audio-only instead of a frozen call."),
    ("googlephotos", "Google Photos", "Hold a big backup until the connection can actually finish it."),
    ("xbox", "Xbox Cloud Gaming", "Block a cloud-gaming session when latency is too high to be playable."),
    ("uber", "Uber", "Confirm the connection before \"Request\", so a tap isn't lost in a dead zone."),
]
ICON = Inches(0.74)
for i, (brand, name, body) in enumerate(apps):
    col = i // 3
    row = i % 3
    x = Inches(0.95 + col * 6.25)
    y = Inches(2.2 + row * 1.6)
    app_icon(s, x, y, ICON, brand)
    text(s, x + Inches(1.02), y - Inches(0.05), Inches(4.85), Inches(1.5),
         [{"text": name, "size": 18, "bold": True, "color": INK, "space_after": 2},
          {"text": body, "size": 13, "color": MUTE, "line_spacing": 1.08}])

# ============================================================ 8. CLOSE
s = slide()
bg(s, INK)
oval(s, Inches(10.6), Inches(4.6), Inches(3.4), Inches(3.4), ORANGE)
kicker(s, Inches(0.9), Inches(1.2), "PingIt SDK", color=ORANGE)
text(s, Inches(0.85), Inches(2.2), Inches(11.2), Inches(2.6),
     [{"text": "A few lines of code. Any app. Users who never hit the wall.",
       "size": 40, "bold": True, "color": CREAM, "line_spacing": 1.08}])
text(s, Inches(0.9), Inches(5.6), Inches(9), Inches(0.5),
     [{"text": "github.com/bitbit138/pingit-sdk", "size": 16, "color": ORANGE, "font": MONO}])

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "PingIt_Pitch.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides)")
